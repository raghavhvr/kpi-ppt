from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def build_ppt(charts, insights_text: str, output_path: str, title: str) -> str:
    prs = Presentation()
    # Title + narrative
    s0 = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    s0.shapes.title.text = title
    s0.placeholders[1].text = insights_text

    # One slide per indicator
    for ind_code, chart_path, caption in charts:
        s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        s.shapes.title.text = caption
        left = Inches(0.5)
        top = Inches(1.6)
        height = Inches(5.0)
        s.shapes.add_picture(str(chart_path), left, top, height=height)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path